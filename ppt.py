from pptx import Presentation

prs = Presentation('E:\\test.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_table:
            continue

        a = [ [ cell.text for cell in row.cells ] for row in shape.table.rows ]
        print(a)